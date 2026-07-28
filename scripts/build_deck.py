#!/usr/bin/env python3
"""Build PRESENTATION.pptx — the ten-slide deck for the 3-5 minute recording.

Ten slides, one accent colour, one grid. Every claim on a slide already appears in
RESPONSE.md; the small diagrams are extracted from that file's Mermaid fences rather
than duplicated here, so the deck cannot drift from the document. The architecture
slide is the exception — it uses assets/architecture.png, drawn separately.

    uv run --with python-pptx scripts/build_deck.py

Two portability notes inherited from scripts/build_pdf.sh:
  - mermaid-cli drives headless Chromium, which fails on Ubuntu 23.10+ with
    "No usable sandbox". Hence the puppeteer config.
  - RESPONSE.md escapes dollar signs so previewers do not parse $...$ as LaTeX
    maths. Anything lifted from it needs them unescaped again.
"""

from __future__ import annotations

import json
import re
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SOURCE = ROOT / "RESPONSE.md"
ARCHITECTURE = ROOT / "assets" / "architecture.png"
OUT = ROOT / "PRESENTATION.pptx"

# --------------------------------------------------------------------------- design

GROUND = RGBColor(0x14, 0x16, 0x1A)
INK = RGBColor(0xF5, 0xF3, 0xEF)
MUTED = RGBColor(0x9A, 0xA0, 0xA8)
FAINT = RGBColor(0x5A, 0x60, 0x68)
ACCENT = RGBColor(0xE0, 0xA9, 0x4A)
BAD = RGBColor(0xC0, 0x56, 0x3E)

# Arial, not something prettier: PowerPoint on another machine will not have a
# designer font, and the substitute it picks is wider — which silently rewraps
# headlines and clips the last line. Arial is present on Windows and macOS, and
# Liberation Sans is metric-identical on Linux, so the layout is what was tested.
SANS = "Arial"
MONO = "Consolas"

# RESPONSE.md's diagrams are painted for a white page; on this ground the green
# subgraph in particular reads as a slab. Darken the fills, keep the strokes so
# the semantic colour coding survives.
DIAGRAM_FILLS = {
    "fill:#0f2f1a": "fill:#161C18",   # the client-account boundary
    "fill:#1a5b2a": "fill:#1B3A24",   # the database
    "fill:#1a3a5b": "fill:#1A2A3A",   # S3
    "fill:#5b1a1a": "fill:#3A1E1C",   # the one-Lambda problem
    "fill:#4a3c1a": "fill:#332B18",   # the MCP adapter
}

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.9)
COL = W - 2 * MARGIN            # 11.53in of usable width
EYEBROW_Y = Inches(0.62)
HEAD_Y = Inches(1.32)
FOOT_Y = Inches(6.82)

# Put your name here — it lands in the footer of every slide.
AUTHOR = "Principal Solutions Architect  ·  take-home review"


# ------------------------------------------------------------------- pptx utilities

def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def write(frame, spans, *, size, color=INK, bold=False, font=SANS,
          spacing=1.15, space_after=0, align=PP_ALIGN.LEFT, first=False):
    """Append a paragraph. `spans` is a string, or (text, {overrides}) tuples."""
    from pptx.oxml.ns import qn  # noqa: PLC0415

    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.alignment = align
    para.line_spacing = spacing
    para.space_after = Pt(space_after)
    # PowerPoint inherits a bullet from the placeholder list style and renders
    # one on every paragraph; LibreOffice does not. Turn it off explicitly.
    props = para._p.get_or_add_pPr()
    props.append(props.makeelement(qn("a:buNone"), {}))
    for span in [spans] if isinstance(spans, str) else spans:
        text, over = (span, {}) if isinstance(span, str) else span
        run = para.add_run()
        run.text = text
        run.font.name = over.get("font", font)
        run.font.size = Pt(over.get("size", size))
        run.font.bold = over.get("bold", bold)
        run.font.color.rgb = over.get("color", color)
    return para


def bar(slide, x, y, w, h, color):
    """A filled rectangle. Connectors would be the obvious choice for a rule,
    but PowerPoint gives them an arrowhead by default — LibreOffice does not,
    so the arrows were invisible until the deck was opened in Office."""
    from pptx.enum.shapes import MSO_SHAPE  # noqa: PLC0415

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def rule(slide, y, *, width=Inches(1.1), color=ACCENT, thickness=Pt(2.5)):
    return bar(slide, MARGIN, y, width, thickness, color)


def slide_base(deck, *, number=None, eyebrow=None, footer=True):
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = GROUND

    if eyebrow:
        frame = textbox(slide, MARGIN, EYEBROW_Y, COL, Inches(0.3))
        label = f"{number:02d}   {eyebrow}" if number else eyebrow
        write(frame, label.upper(), size=12, color=ACCENT, bold=True,
              spacing=1.0, first=True)

    if footer:
        frame = textbox(slide, MARGIN, FOOT_Y, COL, Inches(0.3))
        write(frame, AUTHOR, size=10, color=FAINT, spacing=1.0, first=True)
        if number:
            page = textbox(slide, MARGIN, FOOT_Y, COL, Inches(0.3))
            write(page, f"{number:02d}", size=10, color=FAINT, spacing=1.0,
                  align=PP_ALIGN.RIGHT, first=True)
    return slide


def headline(slide, text, *, size=38, y=HEAD_Y, ruled=True):
    # Half a line of slack per line: if a font substitution rewraps a headline,
    # it grows into the gap instead of being clipped at the box edge.
    lines = text.count("\n") + 1
    height = Inches(0.62 * (lines + 0.5) * size / 38)
    frame = textbox(slide, MARGIN, y, COL, height)
    for i, line in enumerate(text.split("\n")):
        write(frame, line, size=size, bold=True, spacing=1.06, first=(i == 0))
    bottom = y + height + Inches(0.22)
    if ruled:
        rule(slide, bottom)
    return bottom + Inches(0.34)


def notes(slide, seconds, script):
    slide.notes_slide.notes_text_frame.text = f"[~{seconds}s]\n\n{script.strip()}"


def table(slide, y, rows, widths, *, size=15, header=True, height=Inches(0.44)):
    """A grid of text boxes, not a PowerPoint table.

    A real table carries a table *style*, and the default style paints the whole
    graphic frame white with a visible grid. Clearing cell fills and borders is
    not enough — the style also holds the frame at its declared height, so any
    slack below the last row renders as a white slab. Text boxes have no style
    to fight, and lay out identically.
    """
    for r, row in enumerate(rows):
        is_head = header and r == 0
        x = MARGIN
        for c, text in enumerate(row):
            frame = textbox(slide, x, y + height * r, widths[c] - Inches(0.2),
                            height, anchor=MSO_ANCHOR.MIDDLE)
            write(frame, text, size=12 if is_head else size,
                  color=ACCENT if is_head else (INK if c == 0 else MUTED),
                  bold=is_head, spacing=1.12, first=True)
            x += widths[c]
        if is_head:
            bar(slide, MARGIN, y + height - Inches(0.06), Emu(int(sum(widths))),
                Pt(0.75), FAINT)
    return y + height * len(rows)


def callout(slide, y, text, *, size=17, color=INK, height=Inches(0.95)):
    """An indented pull-quote with an accent bar down its left edge."""
    bar(slide, MARGIN, y, Pt(2.5), height, ACCENT)
    frame = textbox(slide, MARGIN + Inches(0.28), y, COL - Inches(0.28), height,
                    anchor=MSO_ANCHOR.MIDDLE)
    for i, line in enumerate(text.split("\n")):
        write(frame, line, size=size, color=color, spacing=1.3, first=(i == 0))
    return y + height


def picture(slide, path, y, *, max_h, max_w=None):
    """Place centred within the column, scaled to fit both bounds."""
    from PIL import Image  # noqa: PLC0415 — optional, only for aspect ratio

    with Image.open(path) as img:
        ratio = img.width / img.height

    max_w = max_w or COL
    height = min(max_h, Emu(int(max_w / ratio)))
    width = Emu(int(height * ratio))
    slide.shapes.add_picture(str(path), MARGIN + (COL - width) // 2, y,
                             width=width, height=height)
    return y + height


# ------------------------------------------------------------------------- diagrams

def render_diagrams(work: Path) -> list[Path]:
    """Render every Mermaid fence in RESPONSE.md to a dark-themed PNG."""
    (work / "puppeteer.json").write_text(
        json.dumps({"args": ["--no-sandbox", "--disable-setuid-sandbox"]}))
    (work / "mermaid.json").write_text(json.dumps({
        "theme": "dark",
        "themeVariables": {
            "background": "#14161A",
            "primaryColor": "#22262C",
            "primaryTextColor": "#F5F3EF",
            "primaryBorderColor": "#4A5058",
            "lineColor": "#9AA0A8",
            "edgeLabelBackground": "#14161A",
            "clusterBkg": "#1B1E23",
            "clusterBorder": "#3A4048",
            "fontFamily": "Lato, sans-serif",
            "fontSize": "17px",
        },
    }))

    blocks = re.findall(r"```mermaid\n(.*?)```", SOURCE.read_text(), re.S)
    out = []
    for i, block in enumerate(blocks, 1):
        mmd, png = work / f"d{i}.mmd", work / f"d{i}.png"
        # The document's inline fills are tuned for a white page. Same shapes,
        # same semantics, darker paint — cosmetic only.
        for light, dark in DIAGRAM_FILLS.items():
            block = block.replace(light, dark)
        mmd.write_text(block)
        subprocess.run(
            ["mmdc", "-p", str(work / "puppeteer.json"),
             "-c", str(work / "mermaid.json"),
             "-i", str(mmd), "-o", str(png),
             "-b", "transparent", "-s", "3", "-w", "1600"],
            check=True, capture_output=True,
        )
        out.append(png)
    return out


# --------------------------------------------------------------------------- slides

def build(diagrams: list[Path]) -> Presentation:
    deck = Presentation()
    deck.slide_width, deck.slide_height = W, H

    # 1 -------------------------------------------------------------- title
    s = slide_base(deck, footer=False)
    frame = textbox(s, MARGIN, Inches(2.5), COL, Inches(0.35))
    write(frame, "PRINCIPAL SOLUTIONS ARCHITECT  ·  TAKE-HOME", size=13,
          color=ACCENT, bold=True, spacing=1.0, first=True)
    frame = textbox(s, MARGIN, Inches(3.0), COL, Inches(1.5))
    write(frame, "Architecture Review", size=60, bold=True, spacing=1.0, first=True)
    rule(s, Inches(4.35), width=Inches(1.6))
    frame = textbox(s, MARGIN, Inches(4.62), COL, Inches(1.0))
    write(frame, "A natural-language agent over legacy real-estate data,",
          size=21, color=MUTED, spacing=1.35, first=True)
    write(frame, "exposed via the Model Context Protocol.", size=21, color=MUTED,
          spacing=1.35)
    frame = textbox(s, MARGIN, Inches(6.55), COL, Inches(0.4))
    write(frame, AUTHOR, size=11, color=FAINT, spacing=1.0, first=True)
    notes(s, 15, """
I was given a junior engineer's draft and asked what's wrong with it and what I'd build
instead. The ingredients are right. What's missing is one boundary, and one distinction
nobody drew. Most of these five minutes go to the three decisions I'd defend hardest.
""")

    # 2 ------------------------------------------------------------ context
    s = slide_base(deck, number=1, eyebrow="Context")
    y = headline(s, "A brokerage wants to query its own data, and act on it", size=32)
    y = table(s, y, [
        ["The constraint", "What it forces on the design"],
        ["~3M records, +5%/month", "5.4M in a year, 9.7M in two. Fitting today is not fitting."],
        ["Raw records never leave the account", "No hosted third parties. Constrains the models — including on the client side."],
        ["Capped monthly budget", "Cost control belongs in the architecture, not a runbook."],
        ["Nightly batch export only", "Freshness capped at 24h, and it must be shown, not hidden."],
        ["Conversational latency", "A few seconds, end to end."],
    ], [Inches(4.0), Inches(7.53)], size=14, height=Inches(0.46))
    callout(s, y + Inches(0.26),
            "The brief never says what format the export is. I read it as a CSV of structured "
            "records — the load-bearing assumption, and the first question I would ask the client.",
            size=15, color=MUTED, height=Inches(0.8))
    notes(s, 27, """
The data lives in an MLS — the cooperative listing database US brokerages run, often
decades old, and the only way out is a nightly file. Five constraints; the underestimated
one is growth, because five percent a month doubles in fourteen.

One assumption up front: the brief never says what format that export is. I read
"records" as a structured CSV, and the document flags what changes if I'm wrong.
""")

    # 3 ------------------------------------------------- the draft, root cause
    s = slide_base(deck, number=2, eyebrow="The draft")
    y = headline(s, "The export file is being used as the database", size=34)
    y = picture(s, diagrams[0], y, max_h=Inches(2.5))
    callout(s, y + Inches(0.35),
            "A file is a delivery mechanism. A database is a query engine. The draft never\n"
            "converts one into the other — so every question pays the conversion again.",
            size=17, height=Inches(0.9))
    notes(s, 28, """
Here's the draft. One Lambda, everything inside it, on every request, from cold: download
the export, scan it in memory, re-embed the corpus, speak MCP, hold the business logic.

Four of the five listed problems are one mistake in different clothes: the file is being
used as the database. A file is a delivery mechanism, a database is a query engine, and
nothing converts one into the other. So every question pays that conversion again.
""")

    # 4 ----------------------------------------------------------- problems
    s = slide_base(deck, number=3, eyebrow="What's wrong")
    y = headline(s, "Six problems. Four are that same mistake.", size=34)
    items = [
        ("The export file is used as the database", "12 seconds and $6.67 per query at 2.2M records", False),
        ("Everything is embedded, indiscriminately", "a price is a fact to compare, not a passage to match", False),
        ("Embeddings recomputed on every query", "~$9 a question for work already done last night", False),
        ("No authentication layer", "an open exfiltration endpoint over 3M confidential records", False),
        ("MCP server and business logic in one function", "one Lambda, one execution role, no seam", False),
        ("No write path, and no freshness signal", "the brief asks the agent to act; it cannot", True),
    ]
    for label, detail, extra in items:
        frame = textbox(s, MARGIN, y, COL, Inches(0.44))
        write(frame, [
            ("→  ", {"color": ACCENT, "bold": True}),
            (label, {"color": INK, "bold": True}),
            ("     " + detail, {"color": MUTED, "size": 14}),
            ("     NOT ON THEIR LIST" if extra else "", {"color": ACCENT, "size": 11, "bold": True}),
        ], size=17, spacing=1.0, first=True)
        y += Inches(0.52)
    notes(s, 22, """
Six problems, graded by severity, and four collapse into that diagnosis. Two were not on
the list I was given: there's no write path at all, though the brief asks the agent to
act, and residency is violated on the client side — the MCP server never calls a model,
the client does.
""")

    # 5 ------------------------------------------------------- architecture
    # Full bleed, no chrome. ARCHITECTURE is 1672x941 — within a thousandth of
    # 16:9, so it lands edge to edge with no letterboxing, and one white slide
    # in a dark deck reads as a deliberate exhibit rather than a pasted image.
    s = slide_base(deck, number=4, footer=False)
    s.shapes.add_picture(str(ARCHITECTURE), 0, 0, width=W, height=H)
    # The diagram's top-left corner is empty white, so the title sits on the
    # image rather than costing it height.
    frame = textbox(s, Inches(0.5), Inches(0.32), Inches(5.5), Inches(0.6))
    write(frame, "Revised architecture", size=26, bold=True,
          color=RGBColor(0x1A, 0x1D, 0x21), spacing=1.0, first=True)
    bar(s, Inches(0.5), Inches(0.92), Inches(1.1), Pt(2.5), ACCENT)
    notes(s, 28, """
The nightly export stays exactly as it is. It just stops pretending to be a database.

It lands in S3, versioned. A nightly job diffs it and updates only what changed — tens of
thousands of rows, not three million. Documents are embedded once, on arrival. Questions
hit an indexed Postgres holding columns and vectors together. The MCP adapter is separate,
with no database permissions of its own, and there's no internet egress at all.
""")

    # 6 --------------------------------------------------------- defend one
    s = slide_base(deck, number=5, eyebrow="Defend #1")
    y = headline(s, "Route by the kind of data, not the kind of question", size=32)
    y = table(s, y, [
        ["Data", "Path", "Rule"],
        ["Price, beds, ZIP, status, dates", "Exact SQL, vetted tools", "Never embedded"],
        ["HOA rules, disclosures, reports", "Vector search over chunks", "Embedded once, at ingest"],
        ["Half-remembered names, addresses", "pg_trgm", "Fuzzy, free, explainable"],
    ], [Inches(4.1), Inches(3.9), Inches(3.53)], size=15, height=Inches(0.55))
    y = callout(s, y + Inches(0.28),
            "Ask for listings under $500,000 and a vector search will hand you a $530,000 property.\n"
            "It is not broken — it is doing exactly what it was built to do. It is the wrong instrument.",
            size=16, height=Inches(0.95))
    # Flows from the callout rather than a fixed y: a wider font substitution
    # rewraps the callout, and a hard coordinate would be overrun by it.
    frame = textbox(s, MARGIN, y + Inches(0.18), COL, Inches(0.4))
    write(frame, "Both live in one PostgreSQL — filter exactly first, then rank semantically "
                 "inside the filtered set.", size=14, color=MUTED, spacing=1.3, first=True)
    notes(s, 45, """
First decision, and the one most likely to be challenged: the answer is not less vector
search or more of it. It's routing.

Structured fields are facts to compare. Embedding a price makes precise information
approximate — ask for under five hundred thousand and you get a five-thirty listing,
because those two descriptions sit close together in meaning. It isn't broken, it's the
wrong instrument — and for a compliance client, subtly wrong is worse than no answer.

Documents are the opposite. "Does this building allow short-term rentals" is buried in a
lawyer's paragraph, never a checkbox. That genuinely needs embeddings — once, on arrival.
Both live in one Postgres because real questions need both halves in one round trip.
""")

    # 7 --------------------------------------------------------- defend two
    s = slide_base(deck, number=6, eyebrow="Defend #2")
    y = headline(s, "Natural language is the interface.\nIt is not the execution engine.", size=34)
    frame = textbox(s, MARGIN, y + Inches(0.1), COL, Inches(0.9))
    write(frame, '"three-bed condos under $500k sitting over 90 days"',
          size=18, color=MUTED, spacing=1.4, first=True)
    write(frame, [("↓", {"color": ACCENT, "bold": True})], size=18, spacing=1.6)
    frame = textbox(s, MARGIN, y + Inches(1.5), COL, Inches(0.5))
    write(frame, "search_listings(bedrooms=3, max_price=500000, min_days_on_market=90)",
          size=17, color=ACCENT, font=MONO, spacing=1.2, first=True)
    callout(s, y + Inches(2.05),
            "The model chooses a pre-approved question and fills in its parameters.\n"
            "It does not write the WHERE clause — a model that writes its own filter can omit\n"
            "the one restricting results to the user's own office.", size=16, height=Inches(1.2))
    notes(s, 25, """
Second decision. The pushback: isn't the point that the AI understands language?

It is. Its job is to turn that sentence into a function call — pick the pre-approved
question, fill in the parameters. Not a vector, and not SQL it wrote itself. A model
writing its own WHERE clause can omit the one restricting results to the user's own
office. That's not a hallucination bug, it's an authorization hole.
""")

    # 8 ------------------------------------------------------- defend three
    s = slide_base(deck, number=7, eyebrow="Defend #3")
    y = headline(s, "Identity travels to the database. On day one.", size=34)
    y = picture(s, diagrams[2], y, max_h=Inches(1.85))
    callout(s, y + Inches(0.32),
            "The model may choose what to ask.\nIt may never choose whose data to ask about.",
            size=19, height=Inches(1.1))
    frame = textbox(s, MARGIN, Inches(5.85), COL, Inches(0.8))
    write(frame, "Authorization is a property of the data layer, not a layer on top — retrofitting it "
                 "means rewriting the data access code you just wrote. And the demo's real audience "
                 "is the compliance team, who can stop this in month three.",
          size=14, color=MUTED, spacing=1.3, first=True)
    notes(s, 38, """
Third decision. "We'll add auth later, once the demo works."

Authorization isn't a layer you put in front, it's a property of the data layer — which
records this person may see becomes a filter inside every query. Adding it later means
rewriting what you just built.

So identity travels the whole way: verified token, principal set server-side, row-level
security in Postgres. A query that forgets to filter still can't return another office's
rows, because the database refuses.

The rule I'd hold under any pressure: the model may choose what to ask. It may never
choose whose data to ask about.
""")

    # 9 ----------------------------------------------------------- measured
    s = slide_base(deck, number=8, eyebrow="Measured")
    y = headline(s, "Both designs, measured on 2.2M real records", size=32)
    y = table(s, y, [
        ["Records", "Draft — data layer", "Revised — data layer"],
        ["100,000", "610 ms  ·  $0.30 / query", "33 ms"],
        ["1,000,000", "5,483 ms  ·  $3.00 / query", "31 ms"],
        ["2,221,849", "11,991 ms  ·  $6.67 / query", "33 ms"],
    ], [Inches(2.6), Inches(5.4), Inches(3.53)], size=16, height=Inches(0.52))
    # Naming the boundary turns the obvious sceptical question — "33 ms is not
    # what the user experiences" — into a point already anticipated.
    y = callout(s, y + Inches(0.26),
                "The right-hand column does not move. An indexed lookup does not care how much data\n"
                "it is not looking at. At 3M records the draft extrapolates to roughly 16 seconds and\n"
                "$9 a question — the revised design's entire monthly floor, spent on one question.",
                size=15, height=Inches(1.15))
    frame = textbox(s, MARGIN, y + Inches(0.18), COL, Inches(0.4))
    write(frame, "Retrieval only — the model adds 1–3 s to both. Which is the point: the database "
                 "is no longer what you wait for.", size=14, color=MUTED, spacing=1.3, first=True)
    notes(s, 33, """
I didn't want to assert any of this, so I built both designs and measured them on two
point two million real records. This is retrieval time — the model adds one to three
seconds to both, so leaving it in would only hide the difference.

The revised column doesn't move — thirty-three milliseconds at ten thousand records and
thirty-three at two million. At three million the draft extrapolates to sixteen seconds
and nine dollars a question: roughly the entire monthly cost of the fixed design, spent
once. The measurement also corrected me — I expected memory to be the wall. It isn't.
""")

    # 10 ------------------------------------------------------------- close
    s = slide_base(deck, number=9, eyebrow="In one line", footer=True)
    frame = textbox(s, MARGIN, Inches(2.05), COL, Inches(1.9))
    write(frame, "Natural language", size=46, bold=True, spacing=1.12, first=True)
    write(frame, "is the interface.", size=46, bold=True, spacing=1.12)
    write(frame, [("It is not the execution engine.", {"color": ACCENT})],
          size=46, bold=True, spacing=1.12)
    rule(s, Inches(4.75), width=Inches(1.6))
    frame = textbox(s, MARGIN, Inches(5.05), COL, Inches(1.2))
    write(frame, "The AI understands what someone means and chooses the right pre-approved question.",
          size=17, color=MUTED, spacing=1.4, first=True)
    write(frame, "The database answers it exactly. Keeping those two jobs separate is what makes this "
                 "fast, affordable, auditable, and signable by compliance.",
          size=17, color=MUTED, spacing=1.4)
    notes(s, 15, """
So what to take from this isn't a list of corrections, it's a way of deciding. Natural
language is the interface. It is not the execution engine. Keeping those two jobs
separate is what makes this fast, affordable, auditable, and signable by compliance.
""")

    return deck


def main() -> None:
    for tool in ("mmdc",):
        if subprocess.run(["which", tool], capture_output=True).returncode:
            sys.exit(f"need {tool}: npm i -g @mermaid-js/mermaid-cli")

    with tempfile.TemporaryDirectory() as tmp:
        work = Path(tmp)
        diagrams = render_diagrams(work)
        print(f"rendered {len(diagrams)} diagrams")
        build(diagrams).save(OUT)

    print(f"wrote {OUT.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
